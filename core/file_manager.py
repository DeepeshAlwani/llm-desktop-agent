import os
import re
import sqlite3
import numpy as np
from datetime import datetime
from pathlib import Path
from typing import Optional
from langchain_ollama import OllamaEmbeddings
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from docx import Document

DB_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "agent_files.db")
WATCHED_FOLDER = os.path.join(os.path.expanduser("~"), "Desktop", "agent_workspace")

def add_formatted_runs(paragraph, text):
    """
    Add formatted runs to an existing paragraph.
    Supports:
        **bold**
        *italic*
        __underline__
        `inline code`
    """

    pattern = r'(\*\*.*?\*\*|__.*?__|\*.*?\*|`.*?`|[^*_`]+)'

    for segment in re.findall(pattern, text):
        run = paragraph.add_run()

        if segment.startswith("**") and segment.endswith("**"):
            run.text = segment[2:-2]
            run.bold = True

        elif segment.startswith("__") and segment.endswith("__"):
            run.text = segment[2:-2]
            run.underline = True

        elif segment.startswith("*") and segment.endswith("*"):
            run.text = segment[1:-1]
            run.italic = True

        elif segment.startswith("`") and segment.endswith("`"):
            run.text = segment[1:-1]
            run.font.name = "Courier New"

        else:
            run.text = segment


def add_formatted_paragraph(doc, line):
    para = doc.add_paragraph()
    add_formatted_runs(para, line)
    return para


def _flush_table(doc, raw_lines):
    """
    Convert markdown table into docx table.
    """

    rows = []

    for line in raw_lines:

        # Skip separator row
        if re.match(r"^\|[\s:\-|]+\|$", line.strip()):
            continue

        row = [
            cell.strip()
            for cell in line.strip().strip("|").split("|")
        ]

        rows.append(row)

    if not rows:
        return

    max_cols = max(len(row) for row in rows)

    table = doc.add_table(rows=len(rows), cols=max_cols)
    table.style = "Table Grid"

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):

            cell = table.cell(r_idx, c_idx)

            # remove default paragraph text
            cell.text = ""

            para = cell.paragraphs[0]
            add_formatted_runs(para, cell_text)


def write_docx(filepath, content):
    document = Document()

    lines = content.split("\n")

    table_buffer = []

    for line in lines:

        stripped = line.strip()

        # ------------------------
        # TABLE HANDLING
        # ------------------------
        if stripped.startswith("|"):
            table_buffer.append(line)
            continue

        if table_buffer:
            _flush_table(document, table_buffer)
            table_buffer = []

        # ------------------------
        # HEADINGS
        # ------------------------
        if line.startswith("#"):

            level = len(line) - len(line.lstrip("#"))

            heading_text = line.lstrip("#").strip()

            # add heading
            para = document.add_heading(level=min(level, 9))

            add_formatted_runs(para, heading_text)

        # ------------------------
        # BULLET LIST
        # ------------------------
        elif stripped.startswith("- ") or stripped.startswith("* "):

            para = document.add_paragraph(style="List Bullet")

            add_formatted_runs(
                para,
                stripped[2:].strip()
            )

        # ------------------------
        # NUMBERED LIST
        # ------------------------
        elif re.match(r"^\d+\.\s+", stripped):

            text = re.sub(
                r"^\d+\.\s+",
                "",
                stripped
            )

            para = document.add_paragraph(style="List Number")

            add_formatted_runs(
                para,
                text
            )

        # ------------------------
        # BLANK LINE
        # ------------------------
        elif stripped == "":
            continue

        # ------------------------
        # NORMAL PARAGRAPH
        # ------------------------
        else:
            add_formatted_paragraph(
                document,
                line
            )

    # Flush final table
    if table_buffer:
        _flush_table(document, table_buffer)

    document.save(filepath)

def _get_conn():
    return sqlite3.connect(DB_PATH)

embedder = OllamaEmbeddings(model="nomic-embed-text-v2-moe")

# supported file types and how to read them
SUPPORTED_EXTENSIONS = {".txt", ".md", ".py", ".js", ".json", ".csv", ".pdf", ".docx", ".xlsx", ".pptx"}

# ---------------------------------------------------------------------------
# Language → file extensions mapping
# ---------------------------------------------------------------------------
LANGUAGE_EXTENSIONS: dict[str, list[str]] = {
    "python":     [".py"],
    "javascript": [".js", ".mjs", ".cjs"],
    "typescript": [".ts", ".tsx"],
    "go":         [".go"],
    "rust":       [".rs"],
    "java":       [".java"],
    "c":          [".c", ".h"],
    "cpp":        [".cpp", ".cc", ".cxx", ".hpp", ".hxx"],
    "ruby":       [".rb"],
    "php":        [".php"],
    "swift":      [".swift"],
    "kotlin":     [".kt", ".kts"],
    "c_sharp":    [".cs"],
    "bash":       [".sh", ".bash"],
    "lua":        [".lua"],
    "scala":      [".scala"],
    "haskell":    [".hs"],
    "elixir":     [".ex", ".exs"],
}

EXT_TO_LANG: dict[str, str] = {
    ext: lang
    for lang, exts in LANGUAGE_EXTENSIONS.items()
    for ext in exts
}

SKIP_DIRS = {
    ".git", ".hg", ".svn",
    "venv", "env", ".venv",
    "node_modules", "__pycache__",
    "migrations", "dist", "build",
    ".mypy_cache", ".pytest_cache",
    "vendor",
}

# ---------------------------------------------------------------------------
# Tree-sitter lazy loader
# ---------------------------------------------------------------------------

_PARSER_CACHE: dict[str, "Parser | None"] = {}  # type: ignore[name-defined]

# Map our internal language keys to the tree-sitter package module names.
# tree-sitter >=0.22 individual packages expose a `language()` function.
_TS_MODULE_NAMES: dict[str, str] = {
    "python":     "tree_sitter_python",
    "javascript": "tree_sitter_javascript",
    "typescript": "tree_sitter_typescript",  # exposes .language_typescript()
    "go":         "tree_sitter_go",
    "rust":       "tree_sitter_rust",
    "java":       "tree_sitter_java",
    "c":          "tree_sitter_c",
    "cpp":        "tree_sitter_cpp",
    "ruby":       "tree_sitter_ruby",
    "php":        "tree_sitter_php",          # exposes .language_php()
    "swift":      "tree_sitter_swift",
    "kotlin":     "tree_sitter_kotlin",
    "c_sharp":    "tree_sitter_c_sharp",
    "bash":       "tree_sitter_bash",
    "lua":        "tree_sitter_lua",
    "scala":      "tree_sitter_scala",
    "haskell":    "tree_sitter_haskell",
    "elixir":     "tree_sitter_elixir",
}

# Some packages expose the binding under a non-default function name.
_TS_LANG_FN: dict[str, str] = {
    "typescript": "language_typescript",
    "php":        "language_php",
}


def _load_parser(language: str):
    """
    Return a configured tree_sitter.Parser, or None if not installed.

    Supports tree-sitter >=0.22 individual packages, e.g.:
        uv pip install tree-sitter tree-sitter-python tree-sitter-javascript ...

    Results are cached so each grammar is only loaded once per process.
    """
    if language in _PARSER_CACHE:
        return _PARSER_CACHE[language]

    module_name = _TS_MODULE_NAMES.get(language)
    if not module_name:
        _PARSER_CACHE[language] = None
        return None

    try:
        import importlib
        from tree_sitter import Language, Parser  # type: ignore

        mod = importlib.import_module(module_name)

        # Resolve the callable that returns the language capsule
        fn_name = _TS_LANG_FN.get(language, "language")
        lang_fn = getattr(mod, fn_name, None)
        if lang_fn is None:
            raise AttributeError(f"{module_name} has no '{fn_name}()'")

        lang_obj = Language(lang_fn())
        parser = Parser(lang_obj)
        _PARSER_CACHE[language] = parser
        return parser

    except Exception as exc:
        # Grammar package not installed — caller will fall back to line-chunks
        _PARSER_CACHE[language] = None
        return None


# ---------------------------------------------------------------------------
# Node-kind queries per language
# "function node kinds" tell us what AST nodes to extract as top-level units
# ---------------------------------------------------------------------------
FUNCTION_NODE_KINDS: dict[str, list[str]] = {
    "python":     ["function_definition", "async_function_definition"],
    "javascript": ["function_declaration", "function_expression",
                   "arrow_function", "method_definition"],
    "typescript": ["function_declaration", "function_expression",
                   "arrow_function", "method_definition", "method_signature"],
    "go":         ["function_declaration", "method_declaration"],
    "rust":       ["function_item", "impl_item"],
    "java":       ["method_declaration", "constructor_declaration"],
    "c":          ["function_definition"],
    "cpp":        ["function_definition"],
    "ruby":       ["method", "singleton_method"],
    "php":        ["function_definition", "method_declaration"],
    "swift":      ["function_declaration", "init_declaration"],
    "kotlin":     ["function_declaration", "anonymous_function"],
    "c_sharp":    ["method_declaration", "constructor_declaration"],
    "bash":       ["function_definition"],
    "lua":        ["function_definition", "local_function"],
    "scala":      ["function_definition", "val_definition"],
    "haskell":    ["function"],
    "elixir":     ["call"],  # def/defp are calls in Elixir's AST
}

NAME_NODE_FIELDS: dict[str, str] = {
    # For most languages the "name" is in a child field called "name"
    "default": "name",
}


def _get_node_name(node, source_bytes: bytes) -> str:
    """Extract the identifier/name from a function/method node."""
    # Try named child "name"
    name_node = node.child_by_field_name("name")
    if name_node:
        return source_bytes[name_node.start_byte:name_node.end_byte].decode("utf-8", errors="replace")

    # For arrow functions and anonymous expressions, use parent context or position
    return f"<anonymous>@L{node.start_point[0]+1}"


def _leading_comment(node, source_bytes: bytes) -> str:
    """
    Grab comment lines / docstrings that appear immediately before a node.
    Works heuristically: look at the preceding sibling or parent's first child.
    """
    parent = node.parent
    if not parent:
        return ""
    siblings = list(parent.children)
    idx = siblings.index(node)
    comments = []
    for i in range(idx - 1, -1, -1):
        sib = siblings[i]
        if sib.type in ("comment", "line_comment", "block_comment",
                        "documentation_comment", "multiline_comment"):
            text = source_bytes[sib.start_byte:sib.end_byte].decode("utf-8", errors="replace")
            comments.insert(0, text.strip())
        elif sib.type in ("newline", "\n", ""):
            continue
        else:
            break
    return "\n".join(comments)


def _python_docstring(node, source_bytes: bytes) -> str:
    """For Python, extract the first expression_statement string inside a function."""
    body = node.child_by_field_name("body")
    if not body:
        return ""
    for child in body.children:
        if child.type == "expression_statement":
            for c in child.children:
                if c.type == "string":
                    raw = source_bytes[c.start_byte:c.end_byte].decode("utf-8", errors="replace")
                    return raw.strip().strip('"""').strip("'''").strip('"').strip("'").strip()
    return ""


# ---------------------------------------------------------------------------
# Tree-sitter based extraction
# ---------------------------------------------------------------------------

def _extract_with_treesitter(
    source: str,
    stored_path: str,      # renamed from filepath, now relative path
    language: str,
) -> list[dict]:
    """Parse `source` with tree-sitter and return a list of function dicts."""
    parser = _load_parser(language)
    if parser is None:
        return []

    source_bytes = source.encode("utf-8")
    try:
        tree = parser.parse(source_bytes)
    except Exception:
        return []

    target_kinds = set(FUNCTION_NODE_KINDS.get(language, []))
    if not target_kinds:
        return []

    results: list[dict] = []

    def walk(node):
        if node.type in target_kinds:
            name = _get_node_name(node, source_bytes)
            raw_source = source_bytes[node.start_byte:node.end_byte].decode("utf-8", errors="replace")

            if language == "python":
                docstring = _python_docstring(node, source_bytes) or _leading_comment(node, source_bytes)
            else:
                docstring = _leading_comment(node, source_bytes)

            start_line = node.start_point[0] + 1
            end_line   = node.end_point[0] + 1

            results.append({
                "name":      name,
                "docstring": docstring,
                "file":      stored_path,          # now the relative path
                "lines":     (start_line, end_line),
                "language":  language,
                "source":    raw_source,
            })

        for child in node.children:
            walk(child)

    walk(tree.root_node)
    return results


# ---------------------------------------------------------------------------
# Fallback: line-chunk extractor (for unsupported file types)
# ---------------------------------------------------------------------------
CHUNK_SIZE = 40  # lines
CHUNK_OVERLAP = 10


def _extract_line_chunks(source: str, stored_path: str, language: str = "text") -> list[dict]:
    """Split the file into overlapping line chunks when no grammar is available."""
    lines = source.splitlines()
    chunks = []
    step = CHUNK_SIZE - CHUNK_OVERLAP
    i = 0
    chunk_idx = 0
    while i < len(lines):
        chunk_lines = lines[i: i + CHUNK_SIZE]
        chunk_text = "\n".join(chunk_lines)
        chunks.append({
            "name":      f"chunk_{chunk_idx}",
            "docstring": "",
            "file":      stored_path,      # relative path
            "lines":     (i + 1, min(i + CHUNK_SIZE, len(lines))),
            "language":  language,
            "source":    chunk_text,
        })
        i += step
        chunk_idx += 1
    return chunks

# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_functions(filepath: str, rel_path: Optional[str] = None) -> list[dict]:
    """
    Parse a single file and return a list of function/method dicts.
    Uses tree-sitter when available, falls back to line-chunking.

    Args:
        filepath: absolute path to the file on disk
        rel_path: path relative to the repository root (stored in metadata)
    """
    ext = Path(filepath).suffix.lower()
    language = EXT_TO_LANG.get(ext)

    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as fh:
            source = fh.read()
    except (OSError, PermissionError):
        return []

    if not source.strip():
        return []

    # If rel_path wasn't provided, fallback to filepath (old behaviour)
    stored_path = rel_path if rel_path is not None else filepath

    if language:
        results = _extract_with_treesitter(source, stored_path, language)
        if results:
            return results
        return _extract_line_chunks(source, stored_path, language)

    return []


def parse_documents(folder: str) -> list[dict]:
    """
    Recursively walk `folder` and extract functions/chunks from all
    recognised source files.
    """
    all_functions: list[dict] = []
    supported_exts = set(EXT_TO_LANG.keys())
    # Normalise to absolute path so relpath works
    folder_abs = os.path.abspath(folder)

    for root, dirs, files in os.walk(folder_abs):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS and not d.startswith(".")]

        for filename in files:
            ext = Path(filename).suffix.lower()
            if ext not in supported_exts:
                continue
            if filename.startswith("__"):
                continue

            full_path = os.path.join(root, filename)
            # Path relative to the repository root
            rel_path = os.path.relpath(full_path, folder_abs)

            # Pass both absolute and relative paths
            funcs = extract_functions(full_path, rel_path)
            if funcs:
                lang = EXT_TO_LANG.get(ext, "unknown")
                print(f"  [{lang}] {len(funcs):3d} units in {full_path}")
            all_functions.extend(funcs)

    return all_functions


def get_file_type(filepath) -> str:
    try:
        file_extension = os.path.splitext(filepath)[1]
        return file_extension
    except Exception as e:
        return f"{e}"
# used to read .txt, .md, .docx
def read_text_files(filepath):
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def read_pdf(filepath):
    import fitz  # pymupdf
    doc = fitz.open(filepath)
    return "\n".join(str(page.get_text()) for page in doc)

def read_csv(filepath):
    import csv
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        return "\n".join(",".join(row) for row in reader)

def read_json(filepath):
    import json
    with open(filepath, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2)

def read_pptx(filepath):
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_val: str = shape.text  # type: ignore[attr-defined]
                if text_val.strip():
                    text.append(text_val)
    return "\n".join(text)

def check_if_pdf_readable(filepath):
    import fitz
    doc = fitz.open(filepath)
    # if first page has no text it's likely scanned
    return bool(doc[0].get_text().strip())

def read_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    return "\n".join(para.text for para in doc.paragraphs)

def read_file_content(filepath):
    if os.path.exists(filepath):
       try:
            file_extension = get_file_type(filepath)
            if file_extension not in SUPPORTED_EXTENSIONS:
                return f"Sorry file not supported"
            
            if file_extension == ".pptx":
                content = read_pptx(filepath)
                return content
            
            elif file_extension == ".pdf":
                if not check_if_pdf_readable(filepath):
                    return "Sorry, we dont support scanned pdf"
                content = read_pdf(filepath)
                return content
            elif file_extension == ".docx":
                content = read_docx(filepath)
            else :
                file_data = read_text_files(filepath)
       except Exception as e:
           return f"{e}"
    else:
        return f"File not found: {filepath}"

def remove_file_from_index(filepath: str):
    conn = _get_conn()
    row = conn.execute(
        "SELECT id FROM documents WHERE filepath = ?", (filepath,)
    ).fetchone()
    if row:
        conn.execute("DELETE FROM chunks WHERE document_id = ?", (row[0],))
        conn.execute("DELETE FROM documents WHERE id = ?", (row[0],))
        conn.commit()
    conn.close()

def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start += chunk_size - overlap  # step back by overlap amount
        return chunks
def index_file(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()
    
    # code files — use tree-sitter for semantic chunking
    if ext in EXT_TO_LANG:
        rel_path = os.path.relpath(filepath, WATCHED_FOLDER)
        chunks_data = extract_functions(filepath, rel_path)
        
        if not chunks_data:
            return f"Could not parse {filepath}"
        
        conn = _get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO documents (filepath, filename, filetype, indexed_at) VALUES (?, ?, ?, ?)",
            (filepath, os.path.basename(filepath), ext, datetime.now().isoformat())
        )
        doc_id = conn.execute(
            "SELECT id FROM documents WHERE filepath = ?", (filepath,)
        ).fetchone()[0]
        conn.execute("DELETE FROM chunks WHERE document_id = ?", (doc_id,))
        
        for i, chunk in enumerate(chunks_data):
            # embed the function name + docstring + source together
            text_to_embed = f"{chunk['name']}\n{chunk['docstring']}\n{chunk['source']}"
            embedding = embedder.embed_query(text_to_embed)
            embedding_bytes = np.array(embedding, dtype=np.float32).tobytes()
            conn.execute(
                "INSERT INTO chunks (document_id, chunk_index, content, embedding) VALUES (?, ?, ?, ?)",
                (doc_id, i, chunk["source"], embedding_bytes)
            )
        
        conn.commit()
        conn.close()
        return f"Indexed {os.path.basename(filepath)} — {len(chunks_data)} functions"
    
    # documents — use read_file_content + character chunking
    else:
        content = read_file_content(filepath)
        if not content or content.startswith("Sorry") or content.startswith("File not"):
            return f"Could not read {filepath}"
        
        chunks = _chunk_text(content, chunk_size=500, overlap=50)
        conn = _get_conn()
        conn.execute(
            "INSERT OR REPLACE INTO documents (filepath, filename, filetype, indexed_at) VALUES (?, ?, ?, ?)",
            (filepath, os.path.basename(filepath), ext, datetime.now().isoformat())
        )
        doc_id = conn.execute(
            "SELECT id FROM documents WHERE filepath = ?", (filepath,)
        ).fetchone()[0]
        conn.execute("DELETE FROM chunks WHERE document_id = ?", (doc_id,))
        
        for i, chunk in enumerate(chunks):
            embedding = embedder.embed_query(chunk)
            embedding_bytes = np.array(embedding, dtype=np.float32).tobytes()
            conn.execute(
                "INSERT INTO chunks (document_id, chunk_index, content, embedding) VALUES (?, ?, ?, ?)",
                (doc_id, i, chunk, embedding_bytes)
            )
        
        conn.commit()
        conn.close()
        return f"Indexed {os.path.basename(filepath)} — {len(chunks)} chunks"

class AgentFileHandler(FileSystemEventHandler):
    def _to_str(self, path) -> str:
        return path.decode() if isinstance(path, bytes) else path

    def on_modified(self, event):
        if not event.is_directory:
            index_file(self._to_str(event.src_path))
    
    def on_created(self, event):
        if not event.is_directory:
            index_file(self._to_str(event.src_path))
    
    def on_deleted(self, event):
        if not event.is_directory:
            remove_file_from_index(self._to_str(event.src_path))

# ---------------------------------------------------------------------------
# Public helper for cosine similarity (used by tools.py)
# ---------------------------------------------------------------------------
def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    if denom == 0:
        return 0.0
    return float(np.dot(a, b) / denom)


def init_file_db():
    """Create the documents and chunks tables if they don't exist."""
    conn = _get_conn()
    conn.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id         INTEGER PRIMARY KEY AUTOINCREMENT,
            filepath   TEXT UNIQUE,
            filename   TEXT,
            filetype   TEXT,
            indexed_at TEXT
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER,
            chunk_index INTEGER,
            content     TEXT,
            embedding   BLOB
        )
    """)
    conn.commit()
    conn.close()