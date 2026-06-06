"""
ppt_renderer.py  —  pure drawing executor for ppt_agent.py.

Responsibilities:
  - Validate contrast of every text/background pair (auto-fixes below 3.0)
  - Execute drawing instructions exactly as given by the agent
  - Fetch/embed images by query string

What this module does NOT do:
  - Choose colors, fonts, layouts, content — the agent decides all of that
  - Apply design defaults beyond what the agent specifies
  - Know anything about palettes, themes, or slide purposes

Never imports ppt_agent or file_manager.
"""

from __future__ import annotations
import io, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colour helpers ────────────────────────────────────────────────────────────

def _rgb(h: str) -> RGBColor:
    try:
        h = h.strip().lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return RGBColor(0xFF, 0xFF, 0xFF)


def _luminance(h: str) -> float:
    try:
        h = h.strip().lstrip("#")
        vals = [int(h[i:i+2], 16) / 255.0 for i in (0, 2, 4)]
        def lin(c):
            return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4
        r, g, b = [lin(v) for v in vals]
        return 0.2126 * r + 0.7152 * g + 0.0722 * b
    except Exception:
        return 1.0


def contrast_ratio(fg: str, bg: str) -> float:
    l1, l2 = _luminance(fg), _luminance(bg)
    lighter, darker = max(l1, l2), min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


# ── Image fetching ────────────────────────────────────────────────────────────

def fetch_image(query: str) -> io.BytesIO | None:
    """
    Fetch an image for the given query using image_search.py (Pixabay + Unsplash).
    Falls back to a solid-color placeholder if no results or keys are missing.
    """
    try:
        from image_search import image_search
        results = image_search(query, max_results=3, fetch_images=True)
        good = [r for r in results if not r.fetch_error and r.base64_data]
        if good:
            import base64
            return io.BytesIO(base64.b64decode(good[0].base64_data))
    except Exception:
        pass

    # Fallback: solid colour placeholder so the slide still renders
    try:
        from PIL import Image
        img = Image.new("RGB", (560, 420), (40, 40, 60))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except Exception:
        return None


# ── Canvas constants ──────────────────────────────────────────────────────────

SW = Inches(13.33)
SH = Inches(7.5)
_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_bg(slide, color_hex: str) -> None:
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = _rgb(color_hex)


# ── Element drawers ───────────────────────────────────────────────────────────

def _draw_rect(slide, el: dict) -> None:
    shape = slide.shapes.add_shape(
        1, Inches(el.get("l", 0)), Inches(el.get("t", 0)),
        Inches(el.get("w", 1)), Inches(el.get("h", 1))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(el["color"])
    if el.get("border_color"):
        shape.line.color.rgb = _rgb(el["border_color"])
    else:
        shape.line.fill.background()


def _draw_text(slide, el: dict) -> None:
    tb = slide.shapes.add_textbox(
        Inches(el.get("l", 0)), Inches(el.get("t", 0)),
        Inches(el.get("w", 6)), Inches(el.get("h", 1))
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = _ALIGN_MAP.get(el.get("align", "left"), PP_ALIGN.LEFT)
    run = p.add_run()
    run.text        = str(el.get("text", ""))
    run.font.size   = Pt(float(el.get("size", 18)))
    run.font.bold   = bool(el.get("bold", False))
    run.font.italic = bool(el.get("italic", False))
    run.font.name   = el.get("font", "Calibri")
    run.font.color.rgb = _rgb(el.get("color", "#FFFFFF"))


def _draw_bullets(slide, el: dict) -> None:
    tb = slide.shapes.add_textbox(
        Inches(el.get("l", 0)), Inches(el.get("t", 0)),
        Inches(el.get("w", 6)), Inches(el.get("h", 5))
    )
    tf = tb.text_frame
    tf.word_wrap = True
    marker = el.get("marker", "▸  ")
    for i, item in enumerate(el.get("items", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(el.get("space_before", 7))
        p.alignment    = PP_ALIGN.LEFT
        run = p.add_run()
        run.text        = marker + str(item)
        run.font.size   = Pt(float(el.get("size", 16)))
        run.font.bold   = bool(el.get("bold", False))
        run.font.italic = bool(el.get("italic", False))
        run.font.name   = el.get("font", "Calibri")
        run.font.color.rgb = _rgb(el.get("color", "#FFFFFF"))


def _draw_image(slide, el: dict) -> None:
    buf = fetch_image(el.get("image_query", "abstract"))
    if buf:
        try:
            slide.shapes.add_picture(
                buf, Inches(el.get("l", 0)), Inches(el.get("t", 0)),
                Inches(el.get("w", 5)), Inches(el.get("h", 4))
            )
            return
        except Exception:
            pass
    _draw_rect(slide, {**el, "color": "#2A3A4A"})   # fallback placeholder


_DRAWERS = {"rect": _draw_rect, "text": _draw_text,
            "bullets": _draw_bullets, "image": _draw_image}


# ── Public API ────────────────────────────────────────────────────────────────

def render_pptx(filepath: str, slides_data: list[dict]) -> dict:
    """
    Render drawing instructions → .pptx file.

    Each slide dict:
    {
      "bg": "#0D1B2A",
      "elements": [
        {"type":"rect",    "l":0,"t":0,"w":13.33,"h":0.4, "color":"#FF941F"},
        {"type":"text",    "l":0.9,"t":1.4,"w":11.2,"h":2.0,
                           "text":"Heading","size":50,"bold":true,
                           "color":"#FFFFFF","align":"left","font":"Trebuchet MS"},
        {"type":"bullets", "l":0.6,"t":2.0,"w":12.0,"h":5.0,
                           "items":["Point one","Point two"],
                           "size":17,"color":"#E8F0E0","font":"Calibri","marker":"▸  "},
        {"type":"image",   "l":7.95,"t":0.9,"w":5.0,"h":6.2,
                           "image_query":"India flag independence day"},
      ]
    }

    Returns:
        {"ok":True,  "path":filepath, "slides":N, "warnings":[...]}
        {"ok":False, "error":"...", "detail":"...", "warnings":[...]}

    Contrast auto-fix: text elements with ratio < threshold get their color
    flipped to #FFFFFF or #111111 automatically; the warning is still logged.
    """
    warnings = []
    try:
        prs = _new_prs()
        os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)

        for slide_data in slides_data:
            slide = _blank(prs)
            bg    = slide_data.get("bg", "#1C1C1C")
            _set_bg(slide, bg)

            for el in slide_data.get("elements", []):
                # ── Contrast check & auto-fix ──────────────────────────────
                if el.get("type") in ("text", "bullets"):
                    fg       = el.get("color", "#FFFFFF")
                    local_bg = el.get("bg_color", bg)
                    cr       = contrast_ratio(fg, local_bg)
                    min_cr   = 3.0 if float(el.get("size", 16)) >= 24 else 4.5
                    if cr < min_cr:
                        snippet = str(el.get("text") or el.get("items", [""])[:1])[:40]
                        warnings.append(
                            f"Contrast fix: '{snippet}' "
                            f"fg={fg} bg={local_bg} ratio={cr:.2f}→auto-fixed"
                        )
                        el["color"] = "#FFFFFF" if _luminance(local_bg) < 0.5 else "#111111"

                # ── Draw ───────────────────────────────────────────────────
                drawer = _DRAWERS.get(el.get("type"))
                if drawer:
                    try:
                        drawer(slide, el)
                    except Exception as e:
                        warnings.append(f"Draw error ({el.get('type')}): {e}")

        prs.save(filepath)
        return {"ok": True, "path": filepath,
                "slides": len(slides_data), "warnings": warnings}

    except Exception as e:
        import traceback
        return {"ok": False, "error": str(e),
                "detail": traceback.format_exc(), "warnings": warnings}